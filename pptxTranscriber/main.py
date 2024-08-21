from pptx import Presentation
import os

# Path to the presentation file
presentation_path = "./lib.pptx"

# Open the presentation
presentation = Presentation(presentation_path)

# Extract text from slides
all_texts = []
for slide in presentation.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            all_texts.append(shape.text)

# Combine all extracted texts
extracted_text = "\n".join(all_texts)

# Get the input file name without extension
input_file_name = os.path.splitext(os.path.basename(presentation_path))[0]

# Create the output directory if it doesn't exist
output_dir = "../out/pptxTranscriber"
os.makedirs(output_dir, exist_ok=True)

# Save the extracted text to a file
output_path = os.path.join(output_dir, f"{input_file_name}.txt")
with open(output_path, "w", encoding="utf-8") as text_file:
    text_file.write(extracted_text)
